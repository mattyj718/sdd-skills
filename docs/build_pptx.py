#!/usr/bin/env python3
"""Build SDD overview PPTX using PowerPoint template."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from copy import deepcopy
import os

TEMPLATE = "template.pptx"
OUTPUT = "sdd-overview.pptx"
DIAGRAMS = "diagrams"

# Brand colors
NAVY = RGBColor(0x06, 0x1A, 0x40)
BRAND_BLUE = RGBColor(0x00, 0x47, 0xD1)
SKY_BLUE = RGBColor(0x21, 0x7D, 0xFE)
LIGHT_BG = RGBColor(0xEB, 0xED, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEAL = RGBColor(0x00, 0xB9, 0xC2)
MAGENTA = RGBColor(0xDE, 0x5D, 0x9F)
YELLOW = RGBColor(0xFF, 0xC5, 0x4F)
GRAY600 = RGBColor(0x6C, 0x78, 0x92)

prs = Presentation(TEMPLATE)

# Delete all existing slides
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

# Layout references
LAYOUT_COVER_DARK = prs.slide_layouts[0]     # 'Cover Dark (1)'
LAYOUT_DIVIDER = prs.slide_layouts[4]         # 'Divider | Dark'
LAYOUT_CONTENT = prs.slide_layouts[6]         # 'Content' - title + subtitle + content
LAYOUT_LIGHT = prs.slide_layouts[30]          # 'Light' - title + subtitle only
LAYOUT_BLANK_LIGHT = prs.slide_layouts[31]    # 'Full Blank | Light'


def set_run(run, text, font_name='Poppins', size=Pt(18), color=NAVY, bold=False, italic=False):
    run.text = text
    run.font.name = font_name
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def add_paragraph(tf, text, font_name='Poppins', size=Pt(18), color=NAVY, bold=False, italic=False, space_before=Pt(4), space_after=Pt(4), alignment=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    run = p.add_run()
    set_run(run, text, font_name, size, color, bold, italic)
    return p


def add_bullet(tf, text, level=0, size=Pt(16), color=NAVY, bold=False):
    p = tf.add_paragraph()
    p.level = level
    p.space_before = Pt(2)
    p.space_after = Pt(2)
    run = p.add_run()
    set_run(run, text, 'Poppins', size, color, bold)
    return p


def add_mixed_paragraph(tf, parts, space_before=Pt(4), space_after=Pt(4)):
    """parts = [(text, color, bold, size), ...]"""
    p = tf.add_paragraph()
    p.space_before = space_before
    p.space_after = space_after
    for text, color, bold, size in parts:
        run = p.add_run()
        set_run(run, text, 'Poppins', size, color, bold)
    return p


def content_slide(title_text, subtitle_text=""):
    """Create a content slide using the Content layout."""
    slide = prs.slides.add_slide(LAYOUT_CONTENT)
    # Set title (idx 0)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = ""
            set_run(ph.text_frame.paragraphs[0].add_run(), title_text, 'Poppins', Pt(28), BRAND_BLUE, True)
        elif ph.placeholder_format.idx == 13:
            ph.text = ""
            if subtitle_text:
                set_run(ph.text_frame.paragraphs[0].add_run(), subtitle_text, 'Poppins', Pt(14), GRAY600)
        elif ph.placeholder_format.idx == 15:
            ph.text = ""  # Clear content placeholder
    return slide


def get_content_tf(slide):
    """Get the content text frame (idx 15) from a content slide."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 15:
            return ph.text_frame
    return None


def light_slide(title_text, subtitle_text=""):
    """Create a light slide (title + subtitle only, no content placeholder)."""
    slide = prs.slides.add_slide(LAYOUT_LIGHT)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = ""
            set_run(ph.text_frame.paragraphs[0].add_run(), title_text, 'Poppins', Pt(28), BRAND_BLUE, True)
        elif ph.placeholder_format.idx == 13:
            ph.text = ""
            if subtitle_text:
                set_run(ph.text_frame.paragraphs[0].add_run(), subtitle_text, 'Poppins', Pt(14), GRAY600)
    return slide


# ============================================================
# SLIDE 1: Title (Cover Dark)
# ============================================================
slide = prs.slides.add_slide(LAYOUT_COVER_DARK)
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0:
        ph.text = ""
        tf = ph.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        set_run(run, "Subagent-driven\ndevelopment", 'Poppins', Pt(44), WHITE, True)
        p2 = tf.add_paragraph()
        p2.space_before = Pt(16)
        run2 = p2.add_run()
        set_run(run2, "A complete multi-agent workflow for turning ideas into shipped code.", 'Poppins', Pt(18), RGBColor(0xB8, 0xDA, 0xFF))

# ============================================================
# SLIDE 2: The 7 skills
# ============================================================
slide = content_slide("The 7 skills")
tf = get_content_tf(slide)
tf.paragraphs[0].text = ""

skills = [
    ("1", "/brainstorming", "Design gate. Explore, question, propose, write spec."),
    ("2", "/writing-plans", "Convert spec into bite-sized TDD implementation plan."),
    ("3", "/subagent-driven-development", "Execute via fresh subagent per task + two-stage review."),
    ("4", "/executing-plans", "Simpler batch execution in a separate session."),
    ("5", "/requesting-code-review", "Dispatch code reviewer with clean context."),
    ("6", "/finishing-a-development-branch", "Verify tests, merge/PR/keep/discard."),
    ("7", "/using-git-worktrees", "Create isolated workspace with safety checks."),
]

from pptx.util import Emu
from pptx.oxml.ns import qn
import lxml.etree as etree

# Build a simple table
rows = len(skills) + 1
cols = 3
# Get content placeholder position
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 15:
        left = ph.left
        top = ph.top
        width = ph.width
        height = ph.height
        break

table_shape = slide.shapes.add_table(rows, cols, left, top, width, Emu(int(height * 0.85)))
table = table_shape.table

# Column widths
table.columns[0].width = Emu(600000)
table.columns[1].width = Emu(4200000)
table.columns[2].width = Emu(6400000)

# Header row
headers = ["#", "Skill", "Purpose"]
for i, h in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    run = p.add_run()
    set_run(run, h, 'Poppins', Pt(12), WHITE, True)
    # Dark blue background
    tcPr = cell._tc.get_or_add_tcPr()
    solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
    srgb = etree.SubElement(solidFill, qn('a:srgbClr'), val='0047D1')

for r, (num, skill, purpose) in enumerate(skills, 1):
    for c, text in enumerate([num, skill, purpose]):
        cell = table.cell(r, c)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        color = BRAND_BLUE if c == 1 else NAVY
        bold = c == 1
        set_run(run, text, 'Poppins', Pt(11), color, bold)
        # Alternating row bg
        if r % 2 == 0:
            tcPr = cell._tc.get_or_add_tcPr()
            solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
            srgb = etree.SubElement(solidFill, qn('a:srgbClr'), val='EBEDF8')

# Remove content placeholder (we used a table instead)
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 15:
        sp = ph._element
        sp.getparent().remove(sp)
        break

# ============================================================
# SLIDE 3: The workflow (diagram)
# ============================================================
slide = light_slide("The workflow")
svg_path = os.path.join(DIAGRAMS, "workflow.svg")
# Add as image - use PNG if available, otherwise convert
png_path = svg_path.replace('.svg', '.png')
if not os.path.exists(png_path):
    import subprocess
    subprocess.run(['mmdc', '-i', svg_path, '-o', png_path, '-b', 'white', '-w', '2400', '-H', '600'],
                   capture_output=True, env={**os.environ, 'NODE_EXTRA_CA_CERTS': '/etc/ssl/cert.pem'})
slide.shapes.add_picture(png_path, Emu(315965), Emu(1800000), Emu(11226830))

# ============================================================
# SLIDE 4: Skill 1 - /brainstorming
# ============================================================
slide = content_slide("Skill 1: /brainstorming", "Hard gate before any implementation. No code without an approved spec.")
tf = get_content_tf(slide)
tf.paragraphs[0].text = ""

steps = [
    "Explores project context (files, docs, recent commits)",
    "Asks clarifying questions (one at a time, multiple choice preferred)",
    "Proposes 2-3 approaches with trade-offs",
    "Writes spec to docs/specs/YYYY-MM-DD-<topic>-design.md",
    "Self-reviews for placeholders, contradictions, scope creep",
    "User reviews the written spec",
    "Transitions to /writing-plans (the only exit)",
]
for i, step in enumerate(steps, 1):
    add_bullet(tf, f"{i}. {step}", size=Pt(14))

add_paragraph(tf, "YAGNI ruthlessly. Design for isolation and clarity.", size=Pt(14), color=BRAND_BLUE, italic=True, space_before=Pt(12))

# ============================================================
# SLIDE 5: Skill 2 - /writing-plans
# ============================================================
slide = content_slide("Skill 2: /writing-plans", "Turn approved specs into construction manuals. Zero codebase context assumed.")
tf = get_content_tf(slide)
tf.paragraphs[0].text = ""

add_paragraph(tf, "Bite-sized steps (2-5 min each), every step is ONE action:", size=Pt(14), color=NAVY, bold=True)
add_paragraph(tf, "", size=Pt(6))

code_lines = [
    "- [ ] Step 1: Write the failing test",
    "      (actual test code here)",
    "- [ ] Step 2: Run it, verify it fails",
    "      Run: pytest tests/test_foo.py -v",
    '      Expected: FAIL with "function not defined"',
    "- [ ] Step 3: Implement minimal code to pass",
    "      (actual implementation code here)",
    "- [ ] Step 4: Run test, verify it passes",
    "- [ ] Step 5: Commit",
]
for line in code_lines:
    add_paragraph(tf, line, font_name='Courier New', size=Pt(11), color=NAVY, space_before=Pt(0), space_after=Pt(0))

# ============================================================
# SLIDE 6: No placeholders rule
# ============================================================
slide = content_slide("Skill 2: No placeholders rule", "These are plan failures. Never write them.")
tf = get_content_tf(slide)
tf.paragraphs[0].text = ""

failures = [
    '"TBD", "TODO", "implement later", "fill in details"',
    '"Add appropriate error handling"',
    '"Write tests for the above" (without actual test code)',
    '"Similar to Task N" (repeat the code)',
    "Steps that describe what to do without showing how",
]
for f in failures:
    add_bullet(tf, f, size=Pt(14))

add_paragraph(tf, "", size=Pt(8))
add_mixed_paragraph(tf, [
    ("Every step has: ", NAVY, False, Pt(14)),
    ("actual code, exact commands, expected output.", BRAND_BLUE, True, Pt(14)),
])

add_paragraph(tf, "Self-review checklist", size=Pt(16), color=BRAND_BLUE, bold=True, space_before=Pt(16))
add_bullet(tf, "1. Spec coverage: can you point to a task for every requirement?", size=Pt(13))
add_bullet(tf, "2. Placeholder scan: any red flag patterns?", size=Pt(13))
add_bullet(tf, "3. Type consistency: do names match across tasks?", size=Pt(13))

# ============================================================
# SLIDE 7: Skill 3 - /subagent-driven-development (diagram)
# ============================================================
slide = light_slide("Skill 3: /subagent-driven-development", "Execute the plan in the current session. Fresh subagent per task. Two-stage review after each.")
svg_path = os.path.join(DIAGRAMS, "task-loop.svg")
png_path = svg_path.replace('.svg', '.png')
if not os.path.exists(png_path):
    import subprocess
    subprocess.run(['mmdc', '-i', svg_path, '-o', png_path, '-b', 'white', '-w', '2400', '-H', '600'],
                   capture_output=True, env={**os.environ, 'NODE_EXTRA_CA_CERTS': '/etc/ssl/cert.pem'})
slide.shapes.add_picture(png_path, Emu(315965), Emu(1800000), Emu(11226830))

# ============================================================
# SLIDE 8: Model selection
# ============================================================
slide = content_slide("Skill 3: Model selection", "Use the least powerful model that can handle each role.")
tf = get_content_tf(slide)
tf.paragraphs[0].text = ""

# Model selection table
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 15:
        left = ph.left
        top = ph.top
        width = ph.width
        break

table_shape = slide.shapes.add_table(4, 3, left, top, width, Emu(1800000))
table = table_shape.table
table.columns[0].width = Emu(3200000)
table.columns[1].width = Emu(4800000)
table.columns[2].width = Emu(3200000)

model_data = [
    ["Task type", "Signal", "Model"],
    ["Mechanical", "1-2 files, clear spec, isolated", "Fast/cheap (Sonnet)"],
    ["Integration", "Multi-file, coordination needed", "Standard (Sonnet)"],
    ["Architecture/Review", "Design judgment, broad context", "Most capable (Opus)"],
]
for r, row in enumerate(model_data):
    for c, text in enumerate(row):
        cell = table.cell(r, c)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        is_header = r == 0
        set_run(run, text, 'Poppins', Pt(12), WHITE if is_header else NAVY, is_header)
        if is_header:
            tcPr = cell._tc.get_or_add_tcPr()
            solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
            srgb = etree.SubElement(solidFill, qn('a:srgbClr'), val='0047D1')
        elif r % 2 == 0:
            tcPr = cell._tc.get_or_add_tcPr()
            solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
            srgb = etree.SubElement(solidFill, qn('a:srgbClr'), val='EBEDF8')

# Status handling below the table
txBox = slide.shapes.add_textbox(left, Emu(top + 2200000), width, Emu(2800000))
tf2 = txBox.text_frame
tf2.word_wrap = True
add_paragraph(tf2, "Implementer status handling", size=Pt(16), color=BRAND_BLUE, bold=True)

statuses = [
    ("DONE", "Proceed to spec review"),
    ("DONE_WITH_CONCERNS", "Read concerns, address if needed"),
    ("NEEDS_CONTEXT", "Provide missing info, re-dispatch"),
    ("BLOCKED", "More context, better model, or smaller pieces"),
]
for status, action in statuses:
    add_mixed_paragraph(tf2, [
        (f"{status}: ", BRAND_BLUE, True, Pt(13)),
        (action, NAVY, False, Pt(13)),
    ], space_before=Pt(2), space_after=Pt(2))

# Remove content placeholder
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 15:
        sp = ph._element
        sp.getparent().remove(sp)
        break

# ============================================================
# SLIDE 9: Skill 4 - /executing-plans
# ============================================================
slide = content_slide("Skill 4: /executing-plans", "Simpler alternative. Separate session, batch mode, periodic review checkpoints.")
tf = get_content_tf(slide)
tf.paragraphs[0].text = ""

add_paragraph(tf, "Process", size=Pt(16), color=BRAND_BLUE, bold=True)
steps = [
    "Load and critically review the plan",
    "Raise concerns before starting",
    "For each task: mark in-progress, follow steps exactly, run verifications",
]
for i, s in enumerate(steps, 1):
    add_bullet(tf, f"{i}. {s}", size=Pt(14))

add_mixed_paragraph(tf, [
    ("4. ", NAVY, False, Pt(14)),
    ("Red/green TDD for every task", BRAND_BLUE, True, Pt(14)),
    (" (tests before implementation, always)", NAVY, False, Pt(14)),
])

add_bullet(tf, "5. When done: /finishing-a-development-branch", size=Pt(14))

add_paragraph(tf, "When to stop", size=Pt(16), color=BRAND_BLUE, bold=True, space_before=Pt(16))
add_bullet(tf, "Hit a blocker (missing dependency, test fails, unclear instruction)", size=Pt(13))
add_bullet(tf, "Plan has critical gaps", size=Pt(13))
add_bullet(tf, "Verification fails repeatedly", size=Pt(13))
add_paragraph(tf, "Ask for clarification rather than guessing.", size=Pt(13), color=GRAY600, italic=True, space_before=Pt(8))

# ============================================================
# SLIDE 10: Skill 5 - /requesting-code-review
# ============================================================
slide = content_slide("Skill 5: /requesting-code-review")
tf = get_content_tf(slide)
tf.paragraphs[0].text = ""

add_paragraph(tf, "When (mandatory)", size=Pt(16), color=BRAND_BLUE, bold=True)
add_bullet(tf, "After each task in subagent-driven development", size=Pt(14))
add_bullet(tf, "After completing a major feature", size=Pt(14))
add_bullet(tf, "Before merge to main", size=Pt(14))

add_paragraph(tf, "How", size=Pt(16), color=BRAND_BLUE, bold=True, space_before=Pt(12))
add_bullet(tf, "1. Get git SHAs (base and head)", size=Pt(14))
add_bullet(tf, "2. Dispatch code-reviewer subagent with: what was built, the plan, the diff", size=Pt(14))

add_paragraph(tf, "Output", size=Pt(16), color=BRAND_BLUE, bold=True, space_before=Pt(12))
add_mixed_paragraph(tf, [("Strengths", BRAND_BLUE, True, Pt(14)), (" (be specific)", NAVY, False, Pt(14))])
add_mixed_paragraph(tf, [("Issues: ", BRAND_BLUE, True, Pt(14)), ("Critical (must fix) / Important (should fix) / Minor (nice to have)", NAVY, False, Pt(14))])
add_mixed_paragraph(tf, [("Assessment: ", BRAND_BLUE, True, Pt(14)), ("Ready to merge? Yes / No / With fixes", NAVY, False, Pt(14))])

# ============================================================
# SLIDE 11: Skill 6 - /finishing-a-development-branch
# ============================================================
slide = content_slide("Skill 6: /finishing-a-development-branch")
tf = get_content_tf(slide)
tf.paragraphs[0].text = ""

add_paragraph(tf, "Process", size=Pt(16), color=BRAND_BLUE, bold=True)
add_mixed_paragraph(tf, [("1. ", NAVY, False, Pt(14)), ("Verify tests pass", BRAND_BLUE, True, Pt(14)), (" (cannot proceed if they fail)", NAVY, False, Pt(14))])
add_mixed_paragraph(tf, [("2. ", NAVY, False, Pt(14)), ("Determine base branch", BRAND_BLUE, True, Pt(14))])
add_mixed_paragraph(tf, [("3. ", NAVY, False, Pt(14)), ("Present exactly 4 options:", BRAND_BLUE, True, Pt(14))])

add_paragraph(tf, "", size=Pt(4))

options = [
    ("1. Merge locally", "Checkout base, pull, merge, test merged result, delete branch"),
    ("2. Push + PR", "Push with -u, gh pr create, display URL"),
    ("3. Keep as-is", "Report status, leave worktree"),
    ("4. Discard", "Requires typed confirmation, deletes branch"),
]
for opt, desc in options:
    add_mixed_paragraph(tf, [
        (f"   {opt}: ", BRAND_BLUE, True, Pt(13)),
        (desc, NAVY, False, Pt(13)),
    ], space_before=Pt(2), space_after=Pt(2))

add_paragraph(tf, "", size=Pt(4))
add_mixed_paragraph(tf, [("4. ", NAVY, False, Pt(14)), ("Cleanup worktree", BRAND_BLUE, True, Pt(14)), (" (for options 1, 2, 4)", NAVY, False, Pt(14))])

# ============================================================
# SLIDE 12: Skill 7 - /using-git-worktrees
# ============================================================
slide = content_slide("Skill 7: /using-git-worktrees", "Isolated workspace with safety verification.")
tf = get_content_tf(slide)
tf.paragraphs[0].text = ""

add_paragraph(tf, "Directory selection", size=Pt(16), color=BRAND_BLUE, bold=True)
add_bullet(tf, "1. Check for .worktrees/ (preferred) or worktrees/", size=Pt(14))
add_bullet(tf, "2. Check CLAUDE.md for preference", size=Pt(14))
add_bullet(tf, "3. Ask user", size=Pt(14))

add_paragraph(tf, "Safety (critical)", size=Pt(16), color=BRAND_BLUE, bold=True, space_before=Pt(12))
add_paragraph(tf, "Verify directory is in .gitignore before creating. If not, add it and commit first.", size=Pt(14))

add_paragraph(tf, "Setup", size=Pt(16), color=BRAND_BLUE, bold=True, space_before=Pt(12))
add_paragraph(tf, "Auto-detects project type and runs the appropriate install command (npm install, pip install, cargo build, go mod download). Then runs baseline tests to verify clean state.", size=Pt(14))

# ============================================================
# SLIDE 13: Design principles
# ============================================================
slide = content_slide("Design principles")
tf = get_content_tf(slide)
tf.paragraphs[0].text = ""

principles = [
    ("Design-first gate:", " no code without an approved spec"),
    ("Two-stage review:", " spec compliance (right thing?) then code quality (built well?)"),
    ("Red/green TDD:", " unconditional. Tests before implementation, always."),
    ("No placeholders:", " every plan step has real code, real commands, expected output"),
    ("Fresh subagents per task:", " prevents context pollution"),
    ("Bite-sized steps:", " 2-5 min granularity, one action per step"),
    ("Isolation via worktrees:", " no branch-switching confusion"),
    ("YAGNI:", " remove unnecessary features ruthlessly"),
]
for i, (label, desc) in enumerate(principles, 1):
    add_mixed_paragraph(tf, [
        (f"{i}. ", NAVY, False, Pt(14)),
        (label, BRAND_BLUE, True, Pt(14)),
        (desc, NAVY, False, Pt(14)),
    ], space_before=Pt(3), space_after=Pt(3))

# ============================================================
# SLIDE 14: Summary
# ============================================================
slide = content_slide("Summary")
tf = get_content_tf(slide)
tf.paragraphs[0].text = ""

add_mixed_paragraph(tf, [
    ("The design stays the ", NAVY, False, Pt(16)),
    ("source of truth.", BRAND_BLUE, True, Pt(16)),
])
add_mixed_paragraph(tf, [
    ("The construction manual is ", NAVY, False, Pt(16)),
    ("disposable", BRAND_BLUE, True, Pt(16)),
    (" and regenerated per milestone.", NAVY, False, Pt(16)),
])

add_paragraph(tf, "For each milestone:", size=Pt(16), color=BRAND_BLUE, bold=True, space_before=Pt(20))
add_paragraph(tf, "/writing-plans  >  /using-git-worktrees  >  /subagent-driven-development  >  /finishing-a-development-branch", font_name='Courier New', size=Pt(12), color=BRAND_BLUE, space_before=Pt(4))

add_paragraph(tf, "What SDD enforces:", size=Pt(16), color=BRAND_BLUE, bold=True, space_before=Pt(20))
enforced = [
    "No code without a spec (brainstorming gate)",
    "No implementation without a plan (writing-plans gate)",
    "Red/green TDD on every task (unconditional)",
    "Two-stage review after every task (spec + quality)",
    "Fresh context per task (no pollution)",
    "Clean finish (tests, merge/PR, worktree cleanup)",
]
for e in enforced:
    add_bullet(tf, e, size=Pt(13))


# ============================================================
# Save
# ============================================================
prs.save(OUTPUT)
print(f"Saved to {OUTPUT}")
